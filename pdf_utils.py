"""
pdf_utils.py — Document Loading Utilities

Supports loading content from: PDF, PPTX, DOCX, TXT, CSV, XLSX files.
Returns a list of LangChain Document objects.
"""

import csv
import openpyxl
from langchain_community.document_loaders import PyPDFLoader
from langchain_core.documents import Document
from pptx import Presentation
from docx import Document as DocxDocument


def load_file(filepath: str, filename: str) -> list:
    """
    Load a document file and return a list of LangChain Document objects.

    Args:
        filepath: Absolute path to the file on disk.
        filename: Original filename (used to determine extension and metadata).

    Returns:
        List of Document objects with page_content and metadata.
    """
    ext = filename.rsplit(".", 1)[-1].lower()
    docs = []

    if ext == "pdf":
        docs = _load_pdf(filepath, filename)
    elif ext == "pptx":
        docs = _load_pptx(filepath, filename)
    elif ext == "docx":
        docs = _load_docx(filepath, filename)
    elif ext == "txt":
        docs = _load_txt(filepath, filename)
    elif ext == "csv":
        docs = _load_csv(filepath, filename)
    elif ext == "xlsx":
        docs = _load_xlsx(filepath, filename)

    return docs


# ── Private Loaders ────────────────────────────────────────────────────────────


def _load_pdf(filepath: str, filename: str) -> list:
    """Load a PDF file using PyPDFLoader."""
    return PyPDFLoader(filepath).load()


def _load_pptx(filepath: str, filename: str) -> list:
    """Load a PowerPoint file slide-by-slide."""
    docs = []
    prs = Presentation(filepath)
    for i, slide in enumerate(prs.slides):
        text = "\n".join(s.text for s in slide.shapes if hasattr(s, "text"))
        if text.strip():
            docs.append(
                Document(
                    page_content=text.strip(),
                    metadata={"source": filename, "slide": i + 1},
                )
            )
    return docs


def _load_docx(filepath: str, filename: str) -> list:
    """Load a Word document."""
    d = DocxDocument(filepath)
    text = "\n".join(p.text for p in d.paragraphs if p.text.strip())
    return [Document(page_content=text, metadata={"source": filename})]


def _load_txt(filepath: str, filename: str) -> list:
    """Load a plain-text file."""
    with open(filepath, "r", errors="ignore") as f:
        return [Document(page_content=f.read(), metadata={"source": filename})]


def _load_csv(filepath: str, filename: str) -> list:
    """Load a CSV file into a single document."""
    with open(filepath, "r", errors="ignore") as f:
        rows = [", ".join(r) for r in csv.reader(f)]
    return [Document(page_content="\n".join(rows), metadata={"source": filename})]


def _load_xlsx(filepath: str, filename: str) -> list:
    """Load an Excel workbook, one document per sheet."""
    docs = []
    wb = openpyxl.load_workbook(filepath)
    for sheet in wb.sheetnames:
        rows = [
            ", ".join(str(c.value) for c in row if c.value)
            for row in wb[sheet].iter_rows()
        ]
        docs.append(
            Document(
                page_content="\n".join(rows),
                metadata={"source": filename, "sheet": sheet},
            )
        )
    return docs
